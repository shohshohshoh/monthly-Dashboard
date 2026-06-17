#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
╔══════════════════════════════════════════════════════════════════════════════╗
║          3D ダッシュボード テンプレート                                      ║
║  カスタマイズ箇所: [CONFIG] タグのある行/ブロックを編集してください           ║
╚══════════════════════════════════════════════════════════════════════════════╝
"""
import os, io, csv, warnings
warnings.filterwarnings("ignore")
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.ticker as mticker
from matplotlib.colors import LinearSegmentedColormap
from mpl_toolkits.mplot3d import Axes3D
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── フォント自動検出 ──────────────────────────────────────────────────────────
import matplotlib.font_manager as fm
_avail = {f.name for f in fm.fontManager.ttflist}
JP = next((f for f in ["Yu Gothic","Meiryo","MS Gothic"] if f in _avail), "sans-serif")
plt.rcParams.update({"font.family": JP, "axes.unicode_minus": False})

# ══════════════════════════════════════════════════════════════════════════════
# [CONFIG] 基本設定
# ══════════════════════════════════════════════════════════════════════════════
TITLE    = "売上分析ダッシュボード"          # スライドタイトル
SUBTITLE = "2023〜2024年度 | 分析対象期間を記入"  # サブタイトル
OUTPUT   = "my_dashboard.pptx"               # 出力ファイル名

# [CONFIG] カラーパレット ─ ここを変えると全体の色調が変わります
COLORS = {
    "BG":    (3,   7,  30),   # スライド背景（最暗部）
    "CARD":  (8,  18,  55),   # カード背景
    "C1":    (0,  180, 255),  # メインアクセント（青系）
    "C2":    (155, 93, 229),  # サブアクセント（紫系）
    "C3":    (0,  245, 160),  # 成功/ポジティブ（緑系）
    "C4":    (255,159,  28),  # 警告/ウォーム（橙系）
    "C5":    (255, 50, 100),  # 危険/ネガティブ（赤系）
    "C6":    (255,214,  10),  # 強調（黄系）
    "C7":    (0,  236, 236),  # 情報（シアン系）
    "GRID":  (15,  45,  74),  # グリッド線
    "TEXT":  (208,232, 255),  # 本文テキスト
    "SUB":   (74, 106, 138),  # サブテキスト
}

# [CONFIG] KPI カード（最大5個）
KPI_ITEMS = [
    {"label": "[KPI名①]", "value": "[値①]", "color": "C1"},
    {"label": "[KPI名②]", "value": "[値②]", "color": "C1"},
    {"label": "[KPI名③]", "value": "[値③]", "color": "C3"},
    {"label": "[KPI名④]", "value": "[値④]", "color": "C2"},
    {"label": "[KPI名⑤]", "value": "[値⑤]", "color": "C2"},
]

# [CONFIG] 外的要因イベント（最大9個）
EVENTS = [
    {"label": "[イベント①]", "color": "C3"},
    {"label": "[イベント②]", "color": "C4"},
    {"label": "[イベント③]", "color": "C6"},
    {"label": "[イベント④]", "color": "C5"},
    {"label": "[イベント⑤]", "color": "C1"},
    {"label": "[イベント⑥]", "color": "C5"},
    {"label": "[イベント⑦]", "color": "C3"},
    {"label": "[イベント⑧]", "color": "C5"},
    {"label": "[イベント⑨]", "color": "C4"},
]

# [CONFIG] 季節性ラベル（12ヶ月分: "超"/"高"/"中"/"低"）
SEASON = ["低","低","高","中","中","高","高","低","中","高","高","超"]

# ══════════════════════════════════════════════════════════════════════════════
# [CONFIG] データ読み込み ─ ここを自分のデータ取得ロジックに置き換えてください
# ══════════════════════════════════════════════════════════════════════════════
DATA_DIR = r".\data"   # データフォルダのパス

def load_data():
    """
    データを読み込んでグローバル変数に格納する関数。
    CSV・Excel・DB など任意のソースに変更可能。

    返り値の例:
        months  : ["1月","2月",...,"12月"]
        series1 : [1000, 1200, ...] （系列1）
        series2 : [1100, 1300, ...] （系列2）
    """
    # ── サンプルデータ（実データに差し替えてください） ──────────────────────
    months  = [f"{i}月" for i in range(1,13)]
    series1 = [1050,870,1280,1150,1020,1380,1450,1230,1320,1540,1660,1920]
    series2 = [1200,980,1450,1320,1180,1540,1620,1380,1490,1710,1830,2100]
    labels_a = ["A社","B社","C社","D社","E社"]
    values_a = [135, 102, 59, 183, 82]
    yoy_a    = [108, 95, 124, 112, 89]
    return months, series1, series2, labels_a, values_a, yoy_a

months, series1, series2, labels_a, values_a, yoy_a = load_data()

# ══════════════════════════════════════════════════════════════════════════════
# 内部ユーティリティ（通常変更不要）
# ══════════════════════════════════════════════════════════════════════════════
def c(key): return COLORS[key]
def rgb(key): return RGBColor(*COLORS[key])

def buf(fig, dpi=160):
    b = io.BytesIO()
    fig.savefig(b, format="png", dpi=dpi, bbox_inches="tight",
                facecolor=fig.get_facecolor(), edgecolor="none")
    b.seek(0); plt.close(fig); return b

def dark_ax(fig, ax, bg=None):
    bg = bg or c("CARD")
    fig.patch.set_facecolor(bg)
    ax.set_facecolor(bg)
    ax.tick_params(colors=c("SUB"), labelsize=7)
    for s in ax.spines.values(): s.set_edgecolor(c("GRID"))
    ax.grid(axis="y", color=c("GRID"), ls="--", lw=0.4, alpha=0.7)

def glow(ax, x, y, col, lw=1.5, label=None):
    for w,a in [(12,.025),(7,.06),(3.5,.15),(lw,1.)]:
        kw=dict(color=col,lw=w,alpha=a,solid_capstyle="round")
        if w==lw and label: kw["label"]=label
        ax.plot(x,y,**kw)

# ══════════════════════════════════════════════════════════════════════════════
# [CONFIG] チャート定義 ─ 各関数を自由に書き換えてください
# ══════════════════════════════════════════════════════════════════════════════

def chart_1():
    """① 3D比較棒グラフ ─ 2系列の月次比較"""
    fig = plt.figure(figsize=(7.2,4.2), facecolor=c("BG"))
    ax  = fig.add_subplot(111, projection="3d")
    ax.set_facecolor(c("BG"))
    for p in [ax.xaxis.pane,ax.yaxis.pane,ax.zaxis.pane]:
        p.fill=False; p.set_edgecolor(c("GRID"))
    ax.xaxis.line.set_color(c("GRID"))
    ax.yaxis.line.set_color(c("GRID"))
    ax.zaxis.line.set_color(c("GRID"))
    for i in range(12):
        ax.bar3d(i,0.05,0,0.36,0.36,series1[i],
                 color=plt.cm.cool(0.2+i/18),alpha=0.88,shade=True)
        ax.bar3d(i,0.60,0,0.36,0.36,series2[i],
                 color=plt.cm.plasma(0.25+i/20),alpha=0.88,shade=True)
    ax.set_xticks(np.arange(12)+0.18)
    ax.set_xticklabels([m.replace("月","") for m in months],fontsize=5.5,color=c("SUB"))
    ax.set_yticks([0.23,0.78])
    ax.set_yticklabels(["系列1","系列2"],fontsize=8,color=c("TEXT"),fontweight="bold")
    ax.tick_params(axis="z",colors=c("SUB"),labelsize=5.5)
    ax.zaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_:f"{int(v):,}"))
    ax.view_init(elev=24,azim=-52)
    ax.set_title("月別3D比較（系列1 vs 系列2）",color=c("TEXT"),fontsize=10,fontweight="bold",pad=6)
    p1=mpatches.Patch(color=plt.cm.cool(0.5),label="系列1")
    p2=mpatches.Patch(color=plt.cm.plasma(0.4),label="系列2")
    ax.legend(handles=[p1,p2],fontsize=7,facecolor=c("CARD"),
              edgecolor=c("GRID"),labelcolor=c("TEXT"),framealpha=0.9)
    fig.tight_layout(pad=0.3)
    return buf(fig,180)

def chart_2():
    """② ネオングロウ折れ線グラフ ─ トレンド分析"""
    fig,ax=plt.subplots(figsize=(5.0,3.2),facecolor=c("BG"))
    ax.set_facecolor(c("BG")); dark_ax(fig,ax,c("BG"))
    x=np.arange(12)
    ax.fill_between(x,series1,alpha=0.07,color=c("C1"))
    ax.fill_between(x,series2,alpha=0.07,color=c("C4"))
    glow(ax,x,series1,c("C1"),label="系列1")
    glow(ax,x,series2,c("C4"),label="系列2")
    ax.scatter(x,series1,color=c("C1"),s=32,zorder=10,edgecolors="w",lw=0.5)
    ax.scatter(x,series2,color=c("C4"),s=32,zorder=10,edgecolors="w",lw=0.5)
    ax.set_xticks(range(12))
    ax.set_xticklabels([m.replace("月","") for m in months],fontsize=6.5,color=c("SUB"))
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_:f"{int(v):,}"))
    ax.set_title("トレンド分析（系列1 vs 系列2）",color=c("TEXT"),fontsize=9,fontweight="bold",pad=5)
    ax.legend(fontsize=7,facecolor=c("CARD"),edgecolor=c("GRID"),labelcolor=c("TEXT"))
    fig.tight_layout(pad=0.4)
    return buf(fig)

def chart_3():
    """③ ドーナツグラフ ─ 構成比"""
    # [CONFIG] ここのデータを差し替えてください
    vals  = [sum(series1)//3, sum(series1)//3, sum(series1)-2*(sum(series1)//3)]
    lbls  = ["カテゴリA","カテゴリB","カテゴリC"]
    cols  = [c("C1"),c("C2"),c("C3")]
    fig,ax=plt.subplots(figsize=(4.0,3.2),facecolor=c("CARD"))
    ax.set_facecolor(c("CARD")); ax.set_aspect("equal")
    _,_,ats=ax.pie(vals,colors=cols,explode=(.05,.05,.05),startangle=90,
                   autopct="%1.1f%%",
                   wedgeprops=dict(width=0.52,edgecolor=c("BG"),linewidth=2.5),
                   pctdistance=0.77)
    for at in ats: at.set_color(c("TEXT")); at.set_fontsize(8); at.set_fontweight("bold")
    circ=plt.Circle((0,0),0.38,color=c("CARD"),zorder=10)
    ax.add_patch(circ)
    ax.text(0,0.07,"合計",ha="center",va="center",fontsize=7.5,color=c("SUB"))
    ax.text(0,-0.1,f"{sum(vals):,}",ha="center",va="center",
            fontsize=9.5,color=c("C1"),fontweight="bold")
    patches=[mpatches.Patch(color=co,label=f"{l}  {v:,}") for co,l,v in zip(cols,lbls,vals)]
    ax.legend(handles=patches,loc="lower center",bbox_to_anchor=(0.5,-0.08),
              ncol=1,fontsize=7,facecolor=c("BG"),edgecolor=c("GRID"),
              labelcolor=c("TEXT"),framealpha=0.9)
    ax.set_title("構成比（ドーナツ）",color=c("TEXT"),fontsize=9,fontweight="bold",pad=5)
    fig.tight_layout(pad=0.4)
    return buf(fig)

def chart_4():
    """④ 2系列比較グロウ折れ線（率・比率向け）"""
    fig,ax=plt.subplots(figsize=(4.5,3.2),facecolor=c("BG"))
    ax.set_facecolor(c("BG")); dark_ax(fig,ax,c("BG"))
    # [CONFIG] ここを差し替えてください（率データ推奨）
    rate1=[v/max(series1)*100 for v in series1]
    rate2=[v/max(series2)*100 for v in series2]
    x=np.arange(12)
    ax.fill_between(x,rate1,alpha=0.10,color=c("C1"))
    ax.fill_between(x,rate2,alpha=0.10,color=c("C2"))
    glow(ax,x,rate1,c("C1"),label="指標1")
    glow(ax,x,rate2,c("C2"),label="指標2")
    ax.scatter(x,rate1,color=c("C1"),s=28,zorder=10,edgecolors="w",lw=0.4)
    ax.scatter(x,rate2,color=c("C2"),s=28,zorder=10,edgecolors="w",lw=0.4)
    ax.set_xticks(range(12))
    ax.set_xticklabels([m.replace("月","") for m in months],fontsize=6.5,color=c("SUB"))
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_:f"{v:.0f}%"))
    ax.set_title("指標比較トレンド",color=c("TEXT"),fontsize=9,fontweight="bold",pad=5)
    ax.legend(fontsize=7,facecolor=c("CARD"),edgecolor=c("GRID"),labelcolor=c("TEXT"))
    fig.tight_layout(pad=0.4)
    return buf(fig)

def chart_5():
    """⑤ グラデーション横棒グラフ（達成率・スコア向け）"""
    fig,ax=plt.subplots(figsize=(4.8,3.2),facecolor=c("CARD"))
    ax.set_facecolor(c("CARD"))
    ax.spines[["top","right","left"]].set_visible(False)
    ax.spines["bottom"].set_edgecolor(c("GRID"))
    ax.grid(axis="x",color=c("GRID"),lw=0.4,alpha=0.6)
    ax.tick_params(colors=c("SUB"),labelsize=7)
    # [CONFIG] ここを差し替えてください
    names=["担当者A","担当者B","担当者C","担当者D","担当者E"]
    scores=[112,101,95,84,71]
    bar_cols=[c("C3") if s>=100 else (c("C6") if s>=90 else c("C5")) for s in scores]
    for i,(s,col) in enumerate(zip(scores,bar_cols)):
        rgb_=np.array(col)/255
        n=120; img=np.zeros((1,n,4))
        img[0,:,0]=rgb_[0]; img[0,:,1]=rgb_[1]; img[0,:,2]=rgb_[2]
        img[0,:,3]=np.linspace(0.15,1.0,n)
        ax.imshow(img,aspect="auto",extent=[0,s,i-0.27,i+0.27],
                  zorder=3,interpolation="bilinear")
        ax.barh(i,s,height=0.55,color="none",edgecolor=col,lw=1.0,zorder=4)
        ax.text(s+0.8,i,f"{s}%",va="center",fontsize=8.5,color=col,fontweight="bold")
    ax.axvline(100,color=c("TEXT"),lw=1,ls="--",alpha=0.45,zorder=5)
    ax.set_xlim(0,125); ax.set_yticks(range(len(names)))
    ax.set_yticklabels(names,fontsize=9,color=c("TEXT"))
    ax.set_xlabel("達成率（%）",fontsize=7,color=c("SUB"))
    ax.set_title("[グラフタイトル]",color=c("TEXT"),fontsize=9,fontweight="bold",pad=5)
    fig.tight_layout(pad=0.4)
    return buf(fig)

def chart_6():
    """⑥ 集合棒グラフ（カテゴリ×系列比較）"""
    fig,ax=plt.subplots(figsize=(4.2,3.2),facecolor=c("CARD"))
    dark_ax(fig,ax,c("CARD"))
    x=np.arange(len(labels_a)); w=0.26
    ax.bar(x-w,values_a,w,color=c("C1"),alpha=0.88,label="系列",zorder=3)
    for i,(lbl,yoy) in enumerate(zip(labels_a,yoy_a)):
        col=c("C3") if yoy>=100 else c("C5")
        ax.text(i,values_a[i]+1.2,f"{yoy}%",ha="center",fontsize=7.5,
                color=col,fontweight="bold")
    ax.set_xticks(x); ax.set_xticklabels(labels_a,fontsize=8,color=c("TEXT"))
    ax.set_ylabel("数量",fontsize=7,color=c("SUB"))
    ax.set_title("[グラフタイトル]",color=c("TEXT"),fontsize=9,fontweight="bold",pad=5)
    ax.legend(fontsize=7,facecolor=c("CARD"),edgecolor=c("GRID"),labelcolor=c("TEXT"))
    fig.tight_layout(pad=0.4)
    return buf(fig)

def chart_7():
    """⑦ ヒートマップ（2系列×12期間）"""
    fig,ax=plt.subplots(figsize=(5.0,2.8),facecolor=c("CARD"))
    ax.set_facecolor(c("CARD"))
    cmap=LinearSegmentedColormap.from_list("neon",
        ["#030B20","#073080","#0070D0","#00B4FF","#00FFD4","#FFD700"],N=256)
    data=np.array([series1,series2],dtype=float)
    im=ax.imshow(data,cmap=cmap,aspect="auto",interpolation="bilinear")
    mn,mx=data.min(),data.max()
    for r in range(2):
        for co in range(12):
            v=data[r,co]; br=(v-mn)/(mx-mn)
            tc=c("TEXT") if br>0.35 else c("SUB")
            ax.text(co,r,f"{int(v):,}",ha="center",va="center",
                    fontsize=6.5,color=tc,fontweight="bold")
    ax.set_xticks(range(12))
    ax.set_xticklabels([m.replace("月","") for m in months],fontsize=7,color=c("SUB"))
    ax.set_yticks([0,1]); ax.set_yticklabels(["系列1","系列2"],fontsize=8,
                                               color=c("TEXT"),fontweight="bold")
    ax.tick_params(length=0)
    for s in ax.spines.values(): s.set_visible(False)
    cb=fig.colorbar(im,ax=ax,fraction=0.022,pad=0.02)
    cb.ax.tick_params(colors=c("SUB"),labelsize=6)
    ax.set_title("[グラフタイトル]",color=c("TEXT"),fontsize=9,fontweight="bold",pad=5)
    fig.tight_layout(pad=0.4)
    return buf(fig)

def chart_8():
    """⑧ 季節性指数棒グラフ（指数・スコア向け）"""
    fig,ax=plt.subplots(figsize=(4.5,3.0),facecolor=c("BG"))
    ax.set_facecolor(c("BG")); dark_ax(fig,ax,c("BG"))
    # [CONFIG] index値を差し替えてください
    avg_m=[(series1[i]+series2[i])/2 for i in range(12)]
    grand=sum(avg_m)/12
    si=[v/grand*100 for v in avg_m]
    bar_cols=[c("C3") if v>=130 else(c("C7") if v>=110 else(c("C6") if v>=90 else c("C5"))) for v in si]
    bars=ax.bar(range(12),si,color=bar_cols,width=0.72,alpha=0.9,edgecolor=c("BG"),lw=0.5,zorder=3)
    for bar,v in zip(bars,si):
        ax.bar(bar.get_x()+0.05,v*0.12,width=bar.get_width()-0.1,
               bottom=v*0.88,color="white",alpha=0.12,zorder=4)
    ax.axhline(100,color=c("TEXT"),lw=1,ls="--",alpha=0.6)
    ax.text(11.5,101.5,"平均",color=c("SUB"),fontsize=6,ha="right")
    ax.set_xticks(range(12))
    ax.set_xticklabels([m.replace("月","") for m in months],fontsize=6.5,color=c("SUB"))
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_:f"{v:.0f}"))
    ax.set_title("[グラフタイトル]",color=c("TEXT"),fontsize=9,fontweight="bold",pad=5)
    leg=[mpatches.Patch(color=c("C3"),label="超ピーク 130+"),
         mpatches.Patch(color=c("C7"),label="高 110〜"),
         mpatches.Patch(color=c("C6"),label="平均 90〜"),
         mpatches.Patch(color=c("C5"),label="低 〜90")]
    ax.legend(handles=leg,fontsize=6,facecolor=c("CARD"),edgecolor=c("GRID"),
              labelcolor=c("TEXT"),ncol=2,loc="lower right")
    fig.tight_layout(pad=0.4)
    return buf(fig)

# ══════════════════════════════════════════════════════════════════════════════
# チャートレンダリング
# ══════════════════════════════════════════════════════════════════════════════
print("チャートを生成中...")
charts = [chart_1(), chart_2(), chart_3(), chart_4(),
          chart_5(), chart_6(), chart_7(), chart_8()]
print(f"  → {len(charts)}チャート完了")

# ══════════════════════════════════════════════════════════════════════════════
# PowerPoint 組み立て（レイアウト固定・通常変更不要）
# ══════════════════════════════════════════════════════════════════════════════
print("PowerPoint を組み立て中...")
prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)
slide = prs.slides.add_slide(prs.slide_layouts[6])

def _rgb(key): return RGBColor(*COLORS[key])

def _shape(l,t,w,h,fill=None,border=None,bw=0.75,g1=None,g2=None,ang=135):
    sh=slide.shapes.add_shape(1,Cm(l),Cm(t),Cm(w),Cm(h))
    if g1 and g2:
        sh.fill.gradient(); sh.fill.gradient_angle=ang
        sh.fill.gradient_stops[0].position=0.0
        sh.fill.gradient_stops[0].color.rgb=RGBColor(*g1)
        sh.fill.gradient_stops[1].position=1.0
        sh.fill.gradient_stops[1].color.rgb=RGBColor(*g2)
    elif fill: sh.fill.solid(); sh.fill.fore_color.rgb=RGBColor(*fill)
    else: sh.fill.background()
    if border: sh.line.color.rgb=RGBColor(*border); sh.line.width=Pt(bw)
    else: sh.line.fill.background()
    return sh

def _txt(text,l,t,w,h,size=9,bold=False,col="TEXT",align=PP_ALIGN.LEFT,italic=False):
    tb=slide.shapes.add_textbox(Cm(l),Cm(t),Cm(w),Cm(h))
    tb.text_frame.word_wrap=True
    p=tb.text_frame.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb=RGBColor(*COLORS[col])

def _img(stream,l,t,w,h):
    slide.shapes.add_picture(stream,Cm(l),Cm(t),Cm(w),Cm(h))

# Background
_shape(0,0,33.87,19.05,g1=COLORS["BG"],g2=(6,14,38),ang=130)
# Title bar
_shape(0,0,33.87,2.2,g1=(4,12,48),g2=(8,22,60),ang=0)
_shape(0,2.18,33.87,0.05,g1=(0,60,150),g2=COLORS["C1"],ang=0)  # glow line
# Title text
_txt(f"◈  {TITLE}",0.4,0.10,22,1.2,size=19,bold=True,col="C1")
_txt(SUBTITLE,0.4,1.35,30,0.75,size=7.5,col="SUB")
# KPI cards
kx0=22.0; kw=2.28; kg=0.12
for i,k in enumerate(KPI_ITEMS):
    lx=kx0+i*(kw+kg)
    _shape(lx,0.12,kw,2.0,g1=(6,18,50),g2=(4,12,38),border=COLORS[k["color"]],bw=0.75)
    _txt(k["label"],lx+0.1,0.18,kw-0.2,0.7,size=6.5,col="SUB",align=PP_ALIGN.CENTER)
    _txt(k["value"],lx+0.1,0.82,kw-0.2,1.0,size=12,bold=True,
         col=k["color"],align=PP_ALIGN.CENTER)
# Events strip
SY=2.28; SH=0.88
_shape(0.2,SY,33.47,SH,g1=(4,10,28),g2=(3,8,22),border=COLORS["GRID"],bw=0.5,ang=0)
_txt("外的要因",0.35,SY+0.05,2.0,0.4,size=7,bold=True,col="C1")
ew=(19.5-2.6)/len(EVENTS)
for i,ev in enumerate(EVENTS):
    ex=2.6+i*ew
    d=slide.shapes.add_shape(9,Cm(ex+ew/2-0.09),Cm(SY+0.33),Cm(0.18),Cm(0.18))
    d.fill.solid(); d.fill.fore_color.rgb=RGBColor(*COLORS[ev["color"]])
    d.line.fill.background()
    _txt(ev["label"],ex+0.02,SY+0.0,ew-0.04,SH,size=5.8,col=ev["color"],align=PP_ALIGN.CENTER)
# Season boxes
scols={"超":(0,200,80),"高":(80,180,60),"中":(200,160,0),"低":(220,60,60)}
sx0=20.3; sbw=(33.47-sx0-0.4)/12
_txt("季節性",sx0-0.05,SY+0.05,1.5,0.35,size=7,bold=True,col="SUB")
for i,(lbl,m) in enumerate(zip(SEASON,months)):
    bx=sx0+i*sbw
    sb=slide.shapes.add_shape(1,Cm(bx),Cm(SY+0.2),Cm(sbw-0.04),Cm(0.6))
    sb.fill.solid(); sb.fill.fore_color.rgb=RGBColor(*scols.get(lbl,(128,128,128)))
    sb.line.fill.background()
    _txt(m.replace("月",""),bx+0.02,SY+0.22,sbw,0.3,
         size=5.3,col="TEXT",align=PP_ALIGN.CENTER)
    _txt(lbl,bx+0.02,SY+0.50,sbw,0.25,
         size=5.5,bold=True,col="TEXT",align=PP_ALIGN.CENTER)
# Charts ── [CONFIG] レイアウト変更はここで
R1=3.26; R2=10.80; RH=7.3
r1=[(0.20,12.50),(12.90,8.50),(21.60,6.20),(28.00,5.87)]
r2=[(0.20,8.30),(8.70,7.30),(16.20,8.50),(24.90,8.77)]
lbls1=["① 月別3D比較","② トレンド分析","③ 構成比","④ 指標比較"]
lbls2=["⑤ 達成率","⑥ カテゴリ別比較","⑦ ヒートマップ","⑧ 指数分析"]
for (lx,lw),im,lb in zip(r1,charts[:4],lbls1):
    _shape(lx,R1,lw,RH,g1=(5,11,30),g2=(3,9,25),border=COLORS["C1"],bw=0.6)
    _img(im,lx,R1,lw,RH)
    _txt(lb,lx+0.15,R1-0.45,lw,0.4,size=7,bold=True,col="C7")
for (lx,lw),im,lb in zip(r2,charts[4:],lbls2):
    _shape(lx,R2,lw,RH,g1=(5,11,30),g2=(3,9,25),border=COLORS["C1"],bw=0.6)
    _img(im,lx,R2,lw,RH)
    _txt(lb,lx+0.15,R2-0.45,lw,0.4,size=7,bold=True,col="C7")
# Footer
_txt("データソース・作成日・備考など",0.3,18.72,33.0,0.33,size=5.5,col="SUB",italic=True)

prs.save(OUTPUT)
print(f"完了: {OUTPUT}")
