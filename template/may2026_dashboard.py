#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
2026年5月 営業日報 ダッシュボード PowerPoint 生成
dashboard_template.py のレイアウトエンジン ×
analysis_report.py   の実データ・チャート
"""
import io, warnings
warnings.filterwarnings("ignore")
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.ticker as mticker
import matplotlib.colors as mcolors
from collections import defaultdict
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import openpyxl

# フォント自動検出
import matplotlib.font_manager as fm
_avail = {f.name for f in fm.fontManager.ttflist}
JP = next((f for f in ["Yu Gothic","Meiryo","MS Gothic"] if f in _avail), "sans-serif")
plt.rcParams.update({"font.family": JP, "axes.unicode_minus": False})

# ══════════════════════════════════════════════════════════════════════════════
# [CONFIG] 基本設定
# ══════════════════════════════════════════════════════════════════════════════
TITLE    = "2026年5月 営業日報 ダッシュボード"
SUBTITLE = "集計期間: 2026年5月1日〜5月31日  ｜  稼働日: 27日 ／ 定休: 4日（月曜）"
OUTPUT   = "my_dashboard_2026年5月.pptx"

# [CONFIG] カラーパレット
COLORS = {
    "BG":   (3,   7,  30),
    "CARD": (8,  18,  55),
    "C1":   (0,  180, 255),
    "C2":   (155, 93, 229),
    "C3":   (0,  245, 160),
    "C4":   (255,159,  28),
    "C5":   (255, 50, 100),
    "C6":   (255,214,  10),
    "C7":   (0,  236, 236),
    "GRID": (15,  45,  74),
    "TEXT": (208,232, 255),
    "SUB":  (74, 106, 138),
}

# [CONFIG] KPI カード（最大5個）― 金額は千円単位
KPI_ITEMS = [
    {"label": "総 売 上 高",  "value": "20,989 千円", "color": "C1"},
    {"label": "純 売 上 高",  "value": "19,097 千円", "color": "C1"},
    {"label": "FOOD 売 上",   "value": "14,975 千円", "color": "C3"},
    {"label": "昼 食 客 数",  "value": "2,520 名",    "color": "C4"},
    {"label": "夕 食 客 数",  "value": "1,534 名",    "color": "C2"},
]

# [CONFIG] トピックス（最大9個）
EVENTS = [
    {"label": "GW連休(5/3-5)",  "color": "C3"},
    {"label": "定休×4日(月)",   "color": "C5"},
    {"label": "宴会コース人気", "color": "C4"},
    {"label": "プレモル首位",   "color": "C1"},
    {"label": "夕食売上>昼食", "color": "C6"},
    {"label": "FOOD比率 71%",   "color": "C3"},
    {"label": "キャッシュレス56%","color":"C7"},
    {"label": "夕食客単価高",   "color": "C2"},
    {"label": "稼働率 87%",     "color": "C4"},
]

# [CONFIG] 季節性ラベル（12ヶ月: 超/高/中/低）
SEASON = ["低","低","中","中","高","超","超","高","中","高","中","低"]

# ══════════════════════════════════════════════════════════════════════════════
# [CONFIG] データ読み込み
# ══════════════════════════════════════════════════════════════════════════════
SRC = r"../data/★営業日報2026年5月.xlsx"

def load_data():
    wb   = openpyxl.load_workbook(SRC, data_only=True)
    ws   = wb["データ"]
    rows = list(ws.iter_rows(min_row=1, max_row=54, values_only=True))

    COL0 = 5
    date_r  = rows[0][COL0:COL0+31]
    hol_r   = rows[1][COL0:COL0+31]
    wday_r  = rows[3][COL0:COL0+31]
    sales_r = rows[4][COL0:COL0+31]   # 純売上高

    days, wdays, closed, daily = [], [], [], []
    for d, h, w, s in zip(date_r, hol_r, wday_r, sales_r):
        if d is not None:
            days.append(d.day)
            wdays.append(w or "")
            closed.append(h == "休")
            daily.append(s or 0)

    ws_i = wb["商品別"]
    food_s, drink_s = defaultdict(int), defaultdict(int)
    for r in ws_i.iter_rows(min_row=3, max_row=95, values_only=True):
        if r[4] and r[7]:   food_s[r[4]]  += (r[7]  or 0)
        if r[8] and r[11]:  drink_s[r[8]] += (r[11] or 0)

    top_food  = sorted(food_s.items(),  key=lambda x: x[1], reverse=True)[:6]
    top_drink = sorted(drink_s.items(), key=lambda x: x[1], reverse=True)[:6]

    M = {
        "total":20_989_657, "cash":8_692_051,
        "jcb":3_806_445,    "chiba":5_643_149,
        "aqua":567_876,     "paypay":1_076_179, "kake":1_215_401,
        "food":14_975_700,  "drink":3_000_230,
        "baiten":230_260,   "sonota":2_783_467,
        "lunch_pax":2_520,  "din_pax":1_534,
        "lunch_amt":9_681_482, "din_amt":11_308_175,
    }
    return days, wdays, closed, daily, top_food, top_drink, M

days, wdays, closed, daily, top_food, top_drink, M = load_data()

# ══════════════════════════════════════════════════════════════════════════════
# 内部ユーティリティ
# ══════════════════════════════════════════════════════════════════════════════
def c(key):    return tuple(v/255 for v in COLORS[key])
def _rgb(key): return RGBColor(*COLORS[key])

def buf(fig, dpi=160):
    b = io.BytesIO()
    fig.savefig(b, format="png", dpi=dpi, bbox_inches="tight",
                facecolor=fig.get_facecolor(), edgecolor="none")
    b.seek(0); plt.close(fig); return b

def dark_ax(ax, fig=None, bg_key="CARD"):
    bg = c(bg_key)
    if fig: fig.patch.set_facecolor(bg)
    ax.set_facecolor(bg)
    ax.tick_params(colors=c("SUB"), labelsize=8)   # +1pt
    for sp in ax.spines.values(): sp.set_edgecolor(c("GRID"))
    ax.grid(axis="y", color=c("GRID"), ls="--", lw=0.4, alpha=0.7)

def glow(ax, x, y, col, lw=1.5, label=None):
    for w, a in [(10,.03),(5,.07),(2.5,.18),(lw,1.)]:
        kw = dict(color=col, lw=w, alpha=a, solid_capstyle="round")
        if w == lw and label: kw["label"] = label
        ax.plot(x, y, **kw)

def y_fmt(ax):
    ax.yaxis.set_major_formatter(
        mticker.FuncFormatter(lambda v,_: f"{int(v/10000)}万"))

def x_fmt(ax):
    ax.xaxis.set_major_formatter(
        mticker.FuncFormatter(lambda v,_: f"{int(v/10000)}万"))

# ══════════════════════════════════════════════════════════════════════════════
# [CONFIG] チャート定義（6枚）
# ══════════════════════════════════════════════════════════════════════════════

def chart_1():
    """① 日次売上トレンド — 棒グラフ＋ネオングロウ折れ線"""
    fig, ax = plt.subplots(figsize=(9.0, 3.8))
    fig.patch.set_facecolor(c("BG"))
    ax.set_facecolor(c("BG"))

    for d, s, cl in zip(days, daily, closed):
        col = c("GRID") if cl else c("C1")
        ax.bar(d, s, color=col, width=0.75,
               alpha=0.35 if cl else 0.5, zorder=2)

    ox = [d for d, cl in zip(days, closed) if not cl]
    oy = [s for s, cl in zip(daily, closed) if not cl]
    glow(ax, ox, oy, c("C1"), lw=1.8, label="純売上高")
    ax.scatter(ox, oy, color=c("C1"), s=20, zorder=10, edgecolors="none")

    avg = np.mean(oy) if oy else 0
    ax.axhline(avg, color=c("C4"), lw=1.2, ls="--", alpha=0.85,
               label=f"稼働日平均 {avg/10000:.0f}万")

    ax.set_xticks(days)
    ax.set_xticklabels(
        [f"{d}\n休" if cl else f"{d}\n{w}"
         for d, w, cl in zip(days, wdays, closed)],
        fontsize=6.2)   # +1pt (5.2→6.2)
    ax.tick_params(colors=c("SUB"))
    y_fmt(ax)
    ax.tick_params(axis="y", colors=c("SUB"), labelsize=7.5)  # +1pt (6.5→7.5)
    for sp in ax.spines.values(): sp.set_edgecolor(c("GRID"))
    ax.grid(axis="y", color=c("GRID"), ls="--", lw=0.4, alpha=0.6)
    ax.set_ylim(0, max(daily) * 1.3)
    ax.set_title("日次売上トレンド（2026年5月）",
                 color=c("TEXT"), fontsize=10, fontweight="bold", pad=5)  # +1pt (9→10)
    ax.legend(fontsize=7.5,                  # +1pt (6.5→7.5)
              facecolor=c("CARD"), edgecolor=c("GRID"),
              labelcolor="white",            # 白
              loc="upper left", framealpha=0.85)
    fig.tight_layout(pad=0.3)
    return buf(fig)

def chart_2():
    """② 曜日別平均売上 — 棒グラフ"""
    order = ["月","火","水","木","金","土","日"]
    wmap  = defaultdict(list)
    for s, w, cl in zip(daily, wdays, closed):
        if not cl and w: wmap[w].append(s)
    avgs  = [np.mean(wmap[w]) if wmap[w] else 0 for w in order]
    cnts  = [len(wmap[w]) for w in order]
    pal   = {
        "月": c("GRID"), "火": c("C1"), "水": c("C1"),
        "木": c("C1"),   "金": c("C4"), "土": c("C3"), "日": c("C3"),
    }

    fig, ax = plt.subplots(figsize=(4.2, 3.5))
    dark_ax(ax, fig)
    bars = ax.bar(order, avgs, color=[pal[w] for w in order],
                  width=0.65, zorder=3, edgecolor=c("BG"), lw=0.5)
    max_a = max(avgs) if avgs else 1
    for bar, avg, n in zip(bars, avgs, cnts):
        if avg > 0:
            ax.text(bar.get_x() + bar.get_width()/2, avg + max_a*0.03,
                    f"{avg/10000:.1f}万\nn={n}",
                    ha="center", va="bottom",
                    fontsize=7.5, fontweight="bold", color=c("TEXT"))  # +1pt (6.5→7.5)
    y_fmt(ax)
    ax.set_title("曜日別 平均売上",
                 color=c("TEXT"), fontsize=10, fontweight="bold", pad=5)  # +1pt (9→10)
    ax.set_ylim(0, max_a * 1.4)
    fig.tight_layout(pad=0.3)
    return buf(fig)

def chart_3():
    """③ 支払方法別 売上構成 — ドーナツグラフ"""
    labels = ["現金","JCB","千葉銀行","アクアコイン","PayPay","売掛金"]
    vals   = [M["cash"],M["jcb"],M["chiba"],M["aqua"],M["paypay"],M["kake"]]
    cols   = [c("C1"),c("C3"),c("C4"),c("C2"),c("C5"),c("C7")]

    fig, ax = plt.subplots(figsize=(4.2, 3.8))
    fig.patch.set_facecolor(c("CARD"))
    ax.set_facecolor(c("CARD"))
    ax.set_aspect("equal")

    _, _, ats = ax.pie(
        vals, colors=cols, startangle=90,
        autopct=lambda p: f"{p:.0f}%" if p > 4 else "",
        wedgeprops=dict(width=0.52, edgecolor=c("BG"), linewidth=2.2),
        pctdistance=0.76, counterclock=False
    )
    for at in ats:
        at.set_color("white")                   # 白
        at.set_fontsize(8)                       # +1pt (7→8)
        at.set_fontweight("bold")

    ax.add_patch(plt.Circle((0,0), 0.38, color=c("CARD"), zorder=10))
    ax.text(0, 0.1,  "総売上高", ha="center", va="center",
            fontsize=8, color=c("SUB"))          # +1pt (7→8)
    ax.text(0,-0.14, f"¥{M['total']:,}", ha="center", va="center",
            fontsize=8.5, color=c("C1"), fontweight="bold")  # +1pt (7.5→8.5)

    patches = [mpatches.Patch(color=co, label=l)
               for co, l in zip(cols, labels)]
    ax.legend(handles=patches, loc="lower center",
              bbox_to_anchor=(0.5,-0.2), ncol=3,
              fontsize=7,                        # +1pt (6→7)
              facecolor=c("BG"), edgecolor=c("GRID"),
              labelcolor="white",                # 白
              framealpha=0.9)
    ax.set_title("支払方法別 構成",
                 color=c("TEXT"), fontsize=10, fontweight="bold", pad=5)  # +1pt (9→10)
    fig.tight_layout(pad=0.3)
    return buf(fig)

def chart_4():
    """④ カテゴリ別 売上金額 — 横棒グラフ"""
    labels = ["FOOD","DRINK","売店","その他"]
    vals   = [M["food"],M["drink"],M["baiten"],M["sonota"]]
    cols   = [c("C4"),c("C1"),c("C3"),c("C2")]
    total  = sum(vals)

    fig, ax = plt.subplots(figsize=(5.0, 3.5))
    dark_ax(ax, fig)
    ax.grid(axis="x", color=c("GRID"), ls="--", lw=0.4, alpha=0.7)
    ax.grid(axis="y", color="none")

    y = np.arange(len(labels))
    ax.barh(y, vals, color=cols, height=0.55,
            edgecolor=c("BG"), lw=0.5, zorder=3)
    ax.set_yticks(y)
    ax.set_yticklabels(labels, fontsize=9.5, color=c("TEXT"))  # +1pt (8.5→9.5)
    max_v = max(vals)
    for i, v in enumerate(vals):
        ax.text(v + max_v*0.025, i,
                f"¥{v/10000:.0f}万  ({v/total*100:.0f}%)",
                va="center", fontsize=8, color=c("TEXT"), fontweight="bold")  # +1pt (7→8)
    x_fmt(ax)
    ax.tick_params(axis="x", colors=c("SUB"), labelsize=7.5)  # +1pt (6.5→7.5)
    ax.set_xlim(0, max_v * 1.45)
    ax.invert_yaxis()
    ax.set_title("カテゴリ別 売上金額",
                 color=c("TEXT"), fontsize=10, fontweight="bold", pad=5)  # +1pt (9→10)
    fig.tight_layout(pad=0.3)
    return buf(fig)

def chart_5():
    """⑤ 昼食 vs 夕食 — 売上・客数・客単価 複合グラフ"""
    labels = ["昼食","夕食"]
    amts   = [M["lunch_amt"], M["din_amt"]]
    pax    = [M["lunch_pax"], M["din_pax"]]
    unit   = [a/p for a,p in zip(amts,pax)]
    cols   = [c("C4"), c("C1")]
    x      = np.arange(2)

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(6.0, 3.5))
    dark_ax(ax1, fig); dark_ax(ax2)
    ax1.grid(axis="x", color="none"); ax2.grid(axis="x", color="none")

    # 左: 売上金額
    bars1 = ax1.bar(x, amts, color=cols, width=0.55, zorder=3,
                    edgecolor=c("BG"), lw=0.5)
    for bar, a in zip(bars1, amts):
        ax1.text(bar.get_x()+bar.get_width()/2, a+max(amts)*0.04,
                 f"{a/10000:.0f}万", ha="center",
                 fontsize=9, fontweight="bold", color=c("TEXT"))  # +1pt (8→9)
    ax1.set_xticks(x)
    ax1.set_xticklabels(labels, fontsize=10, color=c("TEXT"))     # +1pt (9→10)
    y_fmt(ax1)
    ax1.set_ylim(0, max(amts)*1.35)
    ax1.set_title("売上金額",
                  color=c("TEXT"), fontsize=9, fontweight="bold")  # +1pt (8→9)

    # 右: 客数（棒）＋客単価（グロウ折れ線・第2軸）
    ax2r = ax2.twinx()
    ax2r.set_facecolor(c("CARD"))
    ax2.bar(x, pax, color=cols, width=0.55, alpha=0.75,
            zorder=3, edgecolor=c("BG"), lw=0.5)
    glow(ax2r, x, unit, c("C5"), lw=2.0)
    ax2r.scatter(x, unit, color=c("C5"), s=40, zorder=10)
    for xi,(p,u) in enumerate(zip(pax,unit)):
        ax2.text(xi, p+max(pax)*0.05, f"{p:,}名",
                 ha="center", fontsize=9, fontweight="bold", color=c("TEXT"))  # +1pt (8→9)
        ax2r.text(xi, u+max(unit)*0.07, f"¥{u:,.0f}",
                  ha="center", fontsize=8.5, fontweight="bold", color=c("C5"))  # +1pt (7.5→8.5)
    ax2.set_xticks(x)
    ax2.set_xticklabels(labels, fontsize=10, color=c("TEXT"))      # +1pt (9→10)
    ax2.set_ylim(0, max(pax)*1.42)
    ax2r.set_ylim(0, max(unit)*1.5)
    ax2.yaxis.set_major_formatter(
        mticker.FuncFormatter(lambda v,_: f"{int(v):,}"))
    ax2.tick_params(axis="y", colors=c("SUB"), labelsize=7)        # +1pt (6→7)
    ax2r.tick_params(axis="y", colors=c("C5"), labelsize=7)        # +1pt (6→7)
    for sp in ax2r.spines.values(): sp.set_edgecolor(c("GRID"))
    ax2.set_title("客数 & 客単価",
                  color=c("TEXT"), fontsize=9, fontweight="bold")  # +1pt (8→9)

    fig.suptitle("昼食 vs 夕食 比較",
                 color=c("TEXT"), fontsize=10, fontweight="bold")  # +1pt (9→10)
    fig.tight_layout(pad=0.3)
    return buf(fig)

def chart_6():
    """⑥ 売れ筋商品ランキング — FOOD/DRINK Top6 横棒"""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(7.5, 3.8))
    dark_ax(ax1, fig); dark_ax(ax2)

    def _rank(ax, items, base_col, title):
        names = [x[0] for x in items]
        vals  = [x[1] for x in items]
        y     = np.arange(len(names))
        max_v = max(vals) if vals else 1
        grid  = c("GRID")
        t_arr = [0.4 + 0.6 * v/max_v for v in vals]
        bar_cols = [
            tuple(base_col[j]*t + grid[j]*(1-t) for j in range(3))
            for t in t_arr
        ]
        ax.barh(y, vals, color=bar_cols, height=0.6,
                edgecolor=c("BG"), lw=0.5, zorder=3)
        ax.set_yticks(y)
        ax.set_yticklabels(names, fontsize=7.5, color=c("TEXT"))  # +1pt (6.5→7.5)
        ax.invert_yaxis()
        for i, v in enumerate(vals):
            ax.text(v + max_v*0.025, i, f"{v/10000:.0f}万",
                    va="center", fontsize=7.5,               # +1pt (6.5→7.5)
                    color=c("TEXT"), fontweight="bold")
        x_fmt(ax)
        ax.tick_params(axis="x", colors=c("SUB"), labelsize=7)  # +1pt (6→7)
        ax.grid(axis="x", color=c("GRID"), ls="--", lw=0.4, alpha=0.7)
        ax.grid(axis="y", color="none")
        ax.set_xlim(0, max_v * 1.35)
        ax.set_title(title, color=c("TEXT"),
                     fontsize=9.5, fontweight="bold", pad=4)  # +1pt (8.5→9.5)

    _rank(ax1, top_food,  c("C4"), "FOOD ランキング Top6")
    _rank(ax2, top_drink, c("C1"), "DRINK ランキング Top6")
    fig.tight_layout(pad=0.3)
    return buf(fig)

# ── グラフ生成 ──────────────────────────────────
print("チャートを生成中...")
charts = [chart_1(), chart_2(), chart_3(), chart_4(), chart_5(), chart_6()]
print(f"  → {len(charts)}チャート完了")

# ══════════════════════════════════════════════════════════════════════════════
# PowerPoint 組み立て（dashboard_template.py のエンジン流用）
# ══════════════════════════════════════════════════════════════════════════════
print("PowerPoint を組み立て中...")
prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)
slide = prs.slides.add_slide(prs.slide_layouts[6])

def _shape(l,t,w,h,fill=None,border=None,bw=0.75,g1=None,g2=None,ang=135):
    sh = slide.shapes.add_shape(1, Cm(l),Cm(t),Cm(w),Cm(h))
    if g1 and g2:
        sh.fill.gradient(); sh.fill.gradient_angle = ang
        sh.fill.gradient_stops[0].position = 0.0
        sh.fill.gradient_stops[0].color.rgb = RGBColor(*g1)
        sh.fill.gradient_stops[1].position = 1.0
        sh.fill.gradient_stops[1].color.rgb = RGBColor(*g2)
    elif fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(*fill)
    else:
        sh.fill.background()
    if border: sh.line.color.rgb = RGBColor(*border); sh.line.width = Pt(bw)
    else: sh.line.fill.background()
    return sh

def _txt(text,l,t,w,h,size=9,bold=False,col="TEXT",
         align=PP_ALIGN.LEFT,italic=False):
    tb = slide.shapes.add_textbox(Cm(l),Cm(t),Cm(w),Cm(h))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb = RGBColor(*COLORS[col])

def _img(stream,l,t,w,h):
    slide.shapes.add_picture(stream, Cm(l),Cm(t),Cm(w),Cm(h))

# ── 背景・ヘッダー ──────────────────────────────
_shape(0,0,33.87,19.05, g1=COLORS["BG"], g2=(6,14,38), ang=130)
_shape(0,0,33.87,2.20,  g1=(4,12,48),   g2=(8,22,60),  ang=0)
_shape(0,2.18,33.87,0.05, g1=(0,60,150), g2=COLORS["C1"], ang=0)

_txt(f"◈  {TITLE}", 0.4,0.10, 22,1.2, size=18, bold=True,  col="C1")  # +1pt (17→18)
_txt(SUBTITLE,       0.4,1.32, 30,0.75, size=8,  bold=False, col="SUB") # +1pt (7→8)

# ── KPI カード ──────────────────────────────────
kx0=22.0; kw=2.28; kg=0.12
for i, k in enumerate(KPI_ITEMS):
    lx = kx0 + i*(kw+kg)
    _shape(lx,0.12, kw,2.0,
           g1=(6,18,50), g2=(4,12,38), border=COLORS[k["color"]], bw=0.75)
    _txt(k["label"], lx+0.1,0.18, kw-0.2,0.65,
         size=7.0, col="SUB", align=PP_ALIGN.CENTER)   # +1pt (6.0→7.0)
    _txt(k["value"], lx+0.1,0.82, kw-0.2,1.0,
         size=11, bold=True, col=k["color"],            # +1pt (10→11)
         align=PP_ALIGN.CENTER)

# ── トピックス帯 ────────────────────────────────
SY=2.28; SH=0.88
_shape(0.2,SY, 33.47,SH, g1=(4,10,28),g2=(3,8,22), border=COLORS["GRID"],bw=0.5,ang=0)
_txt("トピックス", 0.35,SY+0.05, 2.2,0.4, size=8, bold=True, col="C1")  # +1pt (7→8)
ew = (19.5-2.6) / len(EVENTS)
for i, ev in enumerate(EVENTS):
    ex = 2.6 + i*ew
    d  = slide.shapes.add_shape(9, Cm(ex+ew/2-0.09),Cm(SY+0.33),Cm(0.18),Cm(0.18))
    d.fill.solid(); d.fill.fore_color.rgb = RGBColor(*COLORS[ev["color"]])
    d.line.fill.background()
    _txt(ev["label"], ex+0.02,SY+0.0, ew-0.04,SH,
         size=6.8, col=ev["color"], align=PP_ALIGN.CENTER)  # +1pt (5.8→6.8)

# ── 季節性 ──────────────────────────────────────
months_jp = ["1月","2月","3月","4月","5月","6月",
             "7月","8月","9月","10月","11月","12月"]
scols = {"超":(0,200,80),"高":(80,180,60),"中":(200,160,0),"低":(220,60,60)}
sx0=20.3; sbw=(33.47-sx0-0.4)/12
_txt("季節性", sx0-0.05,SY+0.05, 1.5,0.35, size=8, bold=True, col="SUB")  # +1pt (7→8)
for i,(lbl,m) in enumerate(zip(SEASON, months_jp)):
    bx = sx0 + i*sbw
    sb = slide.shapes.add_shape(1, Cm(bx),Cm(SY+0.2), Cm(sbw-0.04),Cm(0.6))
    sb.fill.solid(); sb.fill.fore_color.rgb = RGBColor(*scols.get(lbl,(128,128,128)))
    sb.line.fill.background()
    _txt(m.replace("月",""), bx+0.02,SY+0.22, sbw,0.3,
         size=6.3, col="TEXT", align=PP_ALIGN.CENTER)   # +1pt (5.3→6.3)
    _txt(lbl, bx+0.02,SY+0.50, sbw,0.25,
         size=6.5, bold=True, col="TEXT", align=PP_ALIGN.CENTER)  # +1pt (5.5→6.5)

# ── チャート配置（3×2 レイアウト）──────────────
R1  = 3.26
R2  = 10.98
RH  = 7.22

row1 = [(0.20,14.70),(15.10,8.60),(23.90, 9.77)]
row2 = [(0.20, 8.90),(9.30,12.40),(21.90,11.77)]

lbls1 = ["① 日次売上トレンド", "③ 支払方法別構成", "④ カテゴリ別売上"]
lbls2 = ["② 曜日別 平均売上",  "⑤ 昼食 vs 夕食", "⑥ 商品ランキング"]

for (lx,lw), im, lb in zip(row1, [charts[0],charts[2],charts[3]], lbls1):
    _shape(lx,R1,lw,RH, g1=(5,11,30),g2=(3,9,25),border=COLORS["C1"],bw=0.6)
    _img(im, lx,R1,lw,RH)
    _txt(lb, lx+0.15,R1-0.44,lw,0.4, size=7.5, bold=True, col="C7")  # +1pt (6.5→7.5)

for (lx,lw), im, lb in zip(row2, [charts[1],charts[4],charts[5]], lbls2):
    _shape(lx,R2,lw,RH, g1=(5,11,30),g2=(3,9,25),border=COLORS["C1"],bw=0.6)
    _img(im, lx,R2,lw,RH)
    _txt(lb, lx+0.15,R2-0.44,lw,0.4, size=7.5, bold=True, col="C7")  # +1pt (6.5→7.5)

# ── フッター ────────────────────────────────────
_txt("データソース: ★営業日報2026年5月.xlsx  ／  Generated by may2026_dashboard.py",
     0.3,18.72, 33.0,0.33, size=6.5, col="SUB", italic=True)  # +1pt (5.5→6.5)

prs.save(OUTPUT)
print(f"完了: {OUTPUT}")
