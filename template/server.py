#!/usr/bin/env python3
"""
ダッシュボード生成 FastAPI バックエンド

起動方法:
  cd template
  uvicorn server:app --port 8000 --reload
"""
import asyncio
import base64
import io
import json
import os
import sys
from collections import defaultdict
from pathlib import Path

import openpyxl
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.table import Table, TableStyleInfo

from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

# パイプライン用スクリプトをモジュールとしてインポート（プロセス起動せず直接呼び出す）
# ※ 重いライブラリ（matplotlib/numpy等）はサーバー起動時に一度だけ読み込まれる
sys.path.insert(0, str(Path(__file__).parent))
import create_daily
import create_dashboard_img
import create_pptx

app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:5173",
        "http://127.0.0.1:5173",
        "https://shohshohshoh.github.io",
    ],
    allow_methods=["*"],
    allow_headers=["*"],
)

ROOT    = Path(__file__).parent.parent
DATA    = ROOT / "frontend" / "public" / "data"
SRC_DIR = ROOT / "data"
SCRIPTS = Path(__file__).parent


class Req(BaseModel):
    year: int
    month: int


# ──────────────────────────────────────────────
# Google Drive ヘルパー
# ──────────────────────────────────────────────

_drive_service_cache = None


def _get_drive_service():
    """Google Drive API サービスを返す（初回のみ構築し、以降はキャッシュを再利用）"""
    global _drive_service_cache
    if _drive_service_cache is not None:
        return _drive_service_cache

    try:
        from google.oauth2 import service_account
        from googleapiclient.discovery import build
    except ImportError:
        raise HTTPException(500, detail="Google Drive SDK が未インストールです")

    creds_json = os.environ.get("GOOGLE_SERVICE_ACCOUNT_JSON")
    if not creds_json:
        raise HTTPException(500, detail="環境変数 GOOGLE_SERVICE_ACCOUNT_JSON が設定されていません")

    credentials = service_account.Credentials.from_service_account_info(
        json.loads(creds_json),
        # 生成物（PPTX/XLSX）を指定フォルダへ自動アップロードするため書き込み権限が必要
        scopes=["https://www.googleapis.com/auth/drive"],
    )
    _drive_service_cache = build("drive", "v3", credentials=credentials)
    return _drive_service_cache


def _download_drive_file(service, file_id: str) -> bytes:
    """Drive からファイルをダウンロードして bytes で返す"""
    from googleapiclient.http import MediaIoBaseDownload
    request = service.files().get_media(fileId=file_id)
    buf = io.BytesIO()
    downloader = MediaIoBaseDownload(buf, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return buf.getvalue()


def _upload_or_replace_file(service, local_path: Path, folder_id: str, mime_type: str) -> str:
    """local_path のファイルを Drive の folder_id 直下へアップロード。
    同名ファイルが既に存在する場合は内容を上書き（再生成のたびに重複が増えないように）。"""
    from googleapiclient.http import MediaFileUpload

    existing = service.files().list(
        q=f"name='{local_path.name}' and '{folder_id}' in parents and trashed=false",
        fields="files(id)",
    ).execute().get("files", [])

    media = MediaFileUpload(str(local_path), mimetype=mime_type, resumable=False)

    if existing:
        file_id = existing[0]["id"]
        service.files().update(fileId=file_id, media_body=media).execute()
        return file_id

    created = service.files().create(
        body={"name": local_path.name, "parents": [folder_id]},
        media_body=media,
        fields="id",
    ).execute()
    return created["id"]


_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


def _sync_outputs_to_drive(y: int, m: int) -> dict:
    """生成済みの pptx・daily xlsx を GOOGLE_DRIVE_UPLOAD_FOLDER_ID へ自動保存する。
    失敗してもメイン処理（生成・ダウンロード提供）は継続させるため例外は投げず結果を dict で返す。"""
    folder_id = (os.environ.get("GOOGLE_DRIVE_UPLOAD_FOLDER_ID") or "").strip()
    if not folder_id:
        return {"uploaded": False, "reason": "GOOGLE_DRIVE_UPLOAD_FOLDER_ID が未設定です"}

    try:
        service = _get_drive_service()
        pptx_path  = DATA / f"dashboard_{y}_{m}.pptx"
        daily_path = DATA / f"daily_{y}_{m}.xlsx"
        _upload_or_replace_file(service, pptx_path,  folder_id, _PPTX_MIME)
        _upload_or_replace_file(service, daily_path, folder_id, _XLSX_MIME)
        return {"uploaded": True}
    except Exception as e:
        return {"uploaded": False, "reason": str(e)}


# ──────────────────────────────────────────────
# パイプライン
# ──────────────────────────────────────────────

def _run_pipeline_sync(y: int, m: int):
    """3ステップを同一プロセス内で直接呼び出す（サブプロセス起動なし）。ブロッキングなので to_thread で呼ぶこと。"""
    steps = [
        ("create_daily",         create_daily.create_daily),
        ("create_dashboard_img", create_dashboard_img.create_dashboard),
        ("create_pptx",          create_pptx.create_pptx),
    ]
    for name, fn in steps:
        try:
            fn(y, m)
        except Exception as e:
            raise HTTPException(500, detail=f"{name} の実行中にエラー: {e}")


async def _run_pipeline(y: int, m: int):
    await asyncio.to_thread(_run_pipeline_sync, y, m)


def _base64_result(y: int, m: int) -> dict:
    png_path   = DATA / f"dashboard_{y}_{m}.png"
    pptx_path  = DATA / f"dashboard_{y}_{m}.pptx"
    daily_path = DATA / f"daily_{y}_{m}.xlsx"
    return {
        "success":        True,
        "png_base64":     base64.b64encode(png_path.read_bytes()).decode(),
        "pptx_base64":    base64.b64encode(pptx_path.read_bytes()).decode(),
        "pptx_filename":  f"dashboard_{y}_{m}.pptx",
        "daily_base64":   base64.b64encode(daily_path.read_bytes()).decode(),
        "daily_filename": f"daily_{y}_{m}.xlsx",
    }


# ──────────────────────────────────────────────
# エンドポイント
# ──────────────────────────────────────────────

@app.post("/api/upload-and-generate")
async def upload_and_generate(
    file: UploadFile = File(...),
    year: int = Form(...),
    month: int = Form(...),
) -> dict:
    y, m = year, month
    SRC_DIR.mkdir(exist_ok=True)
    DATA.mkdir(parents=True, exist_ok=True)
    (SRC_DIR / f"★営業日報{y}年{m}月.xlsx").write_bytes(await file.read())
    await _run_pipeline(y, m)
    result = _base64_result(y, m)
    result["drive_upload"] = await asyncio.to_thread(_sync_outputs_to_drive, y, m)
    return result


@app.post("/api/drive-generate")
async def drive_generate(req: Req) -> dict:
    """Google Drive からソース Excel を取得して生成する"""
    service = _get_drive_service()

    folder_id = (os.environ.get("GOOGLE_DRIVE_FOLDER_ID") or "").strip()
    if not folder_id:
        raise HTTPException(500, detail="環境変数 GOOGLE_DRIVE_FOLDER_ID が設定されていません")

    y, m = req.year, req.month
    filename = f"★営業日報{y}年{m}月.xlsx"
    query = f"name='{filename}' and '{folder_id}' in parents and trashed=false"
    results = service.files().list(q=query, fields="files(id)").execute()
    files = results.get("files", [])
    if not files:
        raise HTTPException(404, detail=f"{filename} が Google Drive フォルダに見つかりません")

    SRC_DIR.mkdir(exist_ok=True)
    DATA.mkdir(parents=True, exist_ok=True)
    (SRC_DIR / filename).write_bytes(_download_drive_file(service, files[0]["id"]))

    await _run_pipeline(y, m)
    result = _base64_result(y, m)
    result["drive_upload"] = await asyncio.to_thread(_sync_outputs_to_drive, y, m)
    return result


@app.get("/api/list-reports")
async def list_reports():
    """Drive 出力フォルダの PPTX ファイルを年月降順で返す"""
    folder_id = (os.environ.get("GOOGLE_DRIVE_OUTPUT_FOLDER_ID") or "").strip()
    if not folder_id:
        return {"reports": []}

    try:
        service = _get_drive_service()
    except HTTPException:
        return {"reports": []}

    results = service.files().list(
        q=f"'{folder_id}' in parents and trashed=false and mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'",
        fields="files(id,name)",
        pageSize=200,
    ).execute()

    reports = []
    for f in results.get("files", []):
        name, fid = f["name"], f["id"]
        # dashboard_YYYY_M.pptx
        if name.startswith("dashboard_") and name.endswith(".pptx"):
            try:
                y, m = name[len("dashboard_"):-len(".pptx")].split("_")
                reports.append({"year": int(y), "month": int(m), "pptx_id": fid})
            except Exception:
                pass

    reports.sort(key=lambda r: (r["year"], r["month"]), reverse=True)
    return {"reports": reports}


@app.get("/api/get-pptx-image/{file_id}")
async def get_pptx_image(file_id: str):
    """Drive の PPTX から最初のスライドの画像を取り出して base64 PNG で返す"""
    from pptx import Presentation

    service = _get_drive_service()
    pptx_bytes = _download_drive_file(service, file_id)

    prs = Presentation(io.BytesIO(pptx_bytes))
    slide = prs.slides[0]
    for shape in slide.shapes:
        if hasattr(shape, "image"):
            return {"png_base64": base64.b64encode(shape.image.blob).decode()}

    raise HTTPException(500, detail="PPTX にスライド画像が見つかりませんでした")


@app.get("/api/get-file/{file_id}")
async def get_file(file_id: str):
    """Drive のファイルを base64 で返す"""
    service = _get_drive_service()
    meta = service.files().get(fileId=file_id, fields="mimeType,name").execute()
    data = _download_drive_file(service, file_id)
    return {
        "base64": base64.b64encode(data).decode(),
        "mime":   meta["mimeType"],
        "name":   meta["name"],
    }


@app.get("/api/combine-daily")
async def combine_daily():
    """Drive 出力フォルダの daily_*.xlsx を全て読み込み、日次データシートを縦結合して返す"""
    folder_id = (os.environ.get("GOOGLE_DRIVE_OUTPUT_FOLDER_ID") or "").strip()
    if not folder_id:
        raise HTTPException(500, detail="GOOGLE_DRIVE_OUTPUT_FOLDER_ID が設定されていません")

    service = _get_drive_service()

    results = service.files().list(
        q=f"'{folder_id}' in parents and trashed=false and mimeType='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'",
        fields="files(id,name)",
        pageSize=200,
    ).execute()

    def _parse_ym(name: str):
        try:
            parts = name[len("daily_"):-len(".xlsx")].split("_")
            return (int(parts[0]), int(parts[1]))
        except Exception:
            return (0, 0)

    files = [f for f in results.get("files", [])
             if f["name"].startswith("daily_") and f["name"].endswith(".xlsx")]
    if not files:
        raise HTTPException(404, detail="daily_*.xlsx ファイルが見つかりません")

    files.sort(key=lambda f: _parse_ym(f["name"]))

    _TABLE_STYLE = TableStyleInfo(
        name="TableStyleMedium9",
        showFirstColumn=False, showLastColumn=False,
        showRowStripes=True, showColumnStripes=False,
    )

    all_rows = []
    headers = None

    for f in files:
        xlsx_bytes = _download_drive_file(service, f["id"])
        wb = openpyxl.load_workbook(io.BytesIO(xlsx_bytes))
        ws = next((wb[s] for s in wb.sheetnames if "日次データ" in s), wb.active)

        row_headers = [ws.cell(1, c).value for c in range(1, ws.max_column + 1)]
        if headers is None:
            headers = row_headers

        for r in range(2, ws.max_row + 1):
            vals = [ws.cell(r, c).value for c in range(1, len(headers) + 1)]
            if vals[0] is None:
                continue
            all_rows.append(vals)

    wb_out = openpyxl.Workbook()
    ws_out = wb_out.active
    ws_out.title = "日次データ（全月）"

    for ci, h in enumerate(headers or [], 1):
        ws_out.cell(1, ci, h)
    for ri, row in enumerate(all_rows, 2):
        for ci, val in enumerate(row, 1):
            cell = ws_out.cell(ri, ci, val)
            if ci == 1 and hasattr(val, "strftime"):
                cell.number_format = "yyyy/mm/dd"

    last_col = get_column_letter(len(headers) if headers else 1)
    last_row = max(1 + len(all_rows), 2)
    tbl = Table(displayName="CombinedDaily", ref=f"A1:{last_col}{last_row}")
    tbl.tableStyleInfo = _TABLE_STYLE
    ws_out.add_table(tbl)

    buf = io.BytesIO()
    wb_out.save(buf)
    return {
        "base64":   base64.b64encode(buf.getvalue()).decode(),
        "filename": "combined_daily.xlsx",
        "count":    len(all_rows),
    }


@app.get("/api/debug-drive")
async def debug_drive():
    """Drive 接続・フォルダアクセス・環境変数の診断"""
    result = {
        "GOOGLE_DRIVE_FOLDER_ID":        os.environ.get("GOOGLE_DRIVE_FOLDER_ID",        "未設定"),
        "GOOGLE_DRIVE_OUTPUT_FOLDER_ID": os.environ.get("GOOGLE_DRIVE_OUTPUT_FOLDER_ID", "未設定"),
        "GOOGLE_DRIVE_UPLOAD_FOLDER_ID": os.environ.get("GOOGLE_DRIVE_UPLOAD_FOLDER_ID", "未設定"),
        "GOOGLE_SERVICE_ACCOUNT_JSON":   "設定あり" if os.environ.get("GOOGLE_SERVICE_ACCOUNT_JSON") else "未設定",
    }
    try:
        service = _get_drive_service()
        result["auth"] = "OK"
    except Exception as e:
        result["auth"] = f"ERROR: {e}"
        return result

    out_id = (os.environ.get("GOOGLE_DRIVE_OUTPUT_FOLDER_ID") or "").strip()
    if out_id:
        try:
            files = service.files().list(
                q=f"'{out_id}' in parents and trashed=false",
                fields="files(name,id)",
                pageSize=20,
            ).execute().get("files", [])
            result["output_folder_accessible"] = True
            result["output_folder_files"] = [f["name"] for f in files]
        except Exception as e:
            result["output_folder_accessible"] = False
            result["output_folder_error"] = str(e)

    up_id = (os.environ.get("GOOGLE_DRIVE_UPLOAD_FOLDER_ID") or "").strip()
    if up_id:
        try:
            meta = service.files().get(
                fileId=up_id,
                fields="id,name,mimeType,driveId,capabilities(canAddChildren,canEdit),owners(emailAddress)",
                supportsAllDrives=True,
            ).execute()
            result["upload_folder_meta"] = meta
        except Exception as e:
            result["upload_folder_meta_error"] = str(e)

        try:
            from googleapiclient.http import MediaInMemoryUpload
            created = service.files().create(
                body={"name": "_upload_test.txt", "parents": [up_id]},
                media_body=MediaInMemoryUpload(b"test", mimetype="text/plain"),
                fields="id,name,parents,driveId",
            ).execute()
            result["upload_test_create"] = created
            service.files().delete(fileId=created["id"], supportsAllDrives=True).execute()
            result["upload_test_cleanup"] = "deleted"
        except Exception as e:
            result["upload_test_create_error"] = str(e)

    return result


@app.get("/api/debug-fonts")
def debug_fonts():
    import glob
    import platform
    import matplotlib.font_manager as _fm

    avail = sorted({f.name for f in _fm.fontManager.ttflist})
    jp_info: dict = {}
    try:
        import japanize_matplotlib as _jm
        pkg = Path(_jm.__file__).parent
        font_files = sorted(pkg.glob("**/*.ttf")) + sorted(pkg.glob("**/*.otf"))
        jp_info = {"loaded": True, "pkg_dir": str(pkg),
                   "font_files": [str(f) for f in font_files]}
    except Exception as e:
        jp_info = {"loaded": False, "error": str(e)}

    sys_fonts = (
        glob.glob("/usr/share/fonts/**/*.ttf", recursive=True) +
        glob.glob("/usr/share/fonts/**/*.otf", recursive=True)
    )
    return {
        "platform": platform.system(),
        "available_fonts_count": len(avail),
        "jp_fonts": [f for f in avail if any(
            k in f for k in ["IPA", "Noto", "Gothic", "Meiryo", "CJK", "BIZ"])],
        "japanize_matplotlib": jp_info,
        "system_font_files": sys_fonts[:20],
    }
