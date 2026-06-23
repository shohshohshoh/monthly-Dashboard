#!/usr/bin/env python3
"""dashboard_yyyy_m.png から PowerPoint スライドを生成する."""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).parent.parent
DATA = ROOT / "frontend" / "public" / "data"


def create_pptx(year: int, month: int) -> Path:
    png_path  = DATA / f"dashboard_{year}_{month}.png"
    pptx_path = DATA / f"dashboard_{year}_{month}.pptx"

    if not png_path.exists():
        raise FileNotFoundError(f"{png_path} が見つかりません")

    prs = Presentation()
    prs.slide_width  = Emu(9144000)   # 10 inch (16:9)
    prs.slide_height = Emu(5143500)   # 5.625 inch

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.add_picture(str(png_path), 0, 0, prs.slide_width, prs.slide_height)

    prs.save(str(pptx_path))
    print(f"保存: {pptx_path}")
    return pptx_path


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("使い方: python create_pptx.py <year> <month>")
        sys.exit(1)
    create_pptx(int(sys.argv[1]), int(sys.argv[2]))
